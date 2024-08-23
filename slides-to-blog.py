from pptx import Presentation
from openai import OpenAI

client = OpenAI(api_key="sk-xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx")
import json

# Ensure you have set up your API key securely, e.g., as an environment variable

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_content = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text.append(paragraph.text)
        if slide_text:
            slides_content.append("\n".join(slide_text))

    return slides_content

def create_blog_post_structure(slides_content):
    blog_post = {
        "title": "Blog Post from Slide Deck",
        "content": []
    }

    for slide_index, content in enumerate(slides_content, start=1):
        lines = content.split('\n')
        if lines:
            heading = lines[0]
            body = "\n".join(lines[1:])
            blog_post["content"].append({
                "heading": f"Slide {slide_index}: {heading}",
                "body": body
            })

    return blog_post

def convert_slides_to_blog_post(slides_content):
    prompt = "Convert the following slide content into a structured blog post. Each slide should be treated as a section in the blog post. Ensure that every section for the blog has atleast 200 words. When writing, write like you are explaining to a new coder or someone just starting to learn technology. Focus on helping folks understand the fundamentals and ensure you do a good job explaining how you move from section of the blog to the next. Ensure the content flows logically from one section to the next.\n\n"
    for i, content in enumerate(slides_content, start=1):
        prompt += f"Slide {i}: {content}\n\n"

    # Use gpt-4 model
    response = client.chat.completions.create(model="gpt-4o",  # Using gpt-4 model
    messages=[
        {"role": "system", "content": "You are a helpful assistant that converts slide content into a blog post."},
        {"role": "user", "content": prompt}
    ],
    max_tokens=1500,  # Adjust based on expected response length
    temperature=0.7)

    blog_post_content = response.choices[0].message.content.strip()
    return blog_post_content

def main(pptx_file_path):
    slides_content = extract_text_from_pptx(pptx_file_path)
    blog_post = create_blog_post_structure(slides_content)

    # Generate the blog post content using AI
    ai_generated_post = convert_slides_to_blog_post(slides_content)

    # Save the AI-generated blog post content
    with open("blog_post.txt", "w") as f:
        f.write(ai_generated_post)

    # Optionally, save the structured data for reference
    with open("blog_post_structure.json", "w") as f:
        json.dump(blog_post, f, indent=2)

if __name__ == "__main__":
    pptx_file_path = "rust.pptx"  # Replace with your file path
    main(pptx_file_path)
