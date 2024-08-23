from pptx import Presentation
from pieces_os_client.wrapper import PiecesClient 
import json


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slides_content = []
    for slide in prs.slides:
        slide_content = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_content += shape.text + "\n"
        slides_content.append(slide_content.strip())
    return slides_content

def create_blog_post_structure(slides_content):
    blog_post = {
        "title": "Blog Post from Presentation",
        "content": []
    }
    for slide_index, content in enumerate(slides_content, start=1):
        lines = content.split('\n')
        if lines:
            heading = lines[0]
            body = "\n".join(lines[1:])
            blog_post["content"].append({
                "heading": f"Slide {slide_index}: {heading}",
                "body": body
            })
    return blog_post

def convert_slides_to_blog_post(slides_content):
    prompt = """You are a helpful assistant that converts slide content into a blog post. 
    Convert the following slide content into a structured blog post using Markdown syntax. 
    Each slide should be treated as a section in the blog post. 
    Ensure that every section for the blog has at least 1000 words. 
    When writing, write like you are explaining to a new coder or someone just starting to learn technology. 
    Focus on helping folks understand the fundamentals and ensure you do a good job explaining how you move from section of the blog to the next. 
    Ensure the content flows logically from one section to the next.
    Use proper Markdown formatting for headers, code blocks, lists, and emphasis where appropriate.

    """
    for i, content in enumerate(slides_content, start=1):
        prompt += f"## Slide {i}\n\n{content}\n\n"

    # Your code to generate content using AI (e.g., GPT-4) goes here
    # For example:
    pieces_client = PiecesClient(config={'baseUrl': 'http://localhost:1000'})  # Assuming you have this set up
    response = pieces_client.copilot.question(prompt).answers.iterable[0].text

    return response  # This should now be properly formatted Markdown

def main(pptx_file_path):
    slides_content = extract_text_from_pptx(pptx_file_path)
    blog_post = create_blog_post_structure(slides_content)

    # Generate the blog post content using AI
    ai_generated_post = convert_slides_to_blog_post(slides_content)

    # Save the AI-generated blog post content as Markdown
    with open("blog_post.md", "w") as f:
        f.write(ai_generated_post)

    # Optionally, save the structured data for reference
    with open("blog_post_structure.json", "w") as f:
        json.dump(blog_post, f, indent=2)

if __name__ == "__main__":
    pptx_file_path = "rust.pptx"  # Replace with your file path
    main(pptx_file_path)
