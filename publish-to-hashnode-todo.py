import streamlit as st
from pptx import Presentation
from pieces_os_client.wrapper import PiecesClient 
import json
import requests

def extract_text_from_pptx(file):
    prs = Presentation(file)
    slides_content = []
    for slide in prs.slides:
        slide_content = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_content += shape.text + "\n"
        slides_content.append(slide_content.strip())
    return slides_content

def create_blog_post_structure(slides_content):
    blog_post = {
        "title": "Blog Post from Presentation",
        "content": []
    }
    for slide_index, content in enumerate(slides_content, start=1):
        lines = content.split('\n')
        if lines:
            heading = lines[0]
            body = "\n".join(lines[1:])
            blog_post["content"].append({
                "heading": f"Slide {slide_index}: {heading}",
                "body": body
            })
    return blog_post

# def convert_slides_to_blog_post(slides_content):
#     prompt = """You are a helpful assistant that converts slide content into a blog post. 
#     Convert the following slide content into a structured blog post using Markdown syntax. 
#     Each slide should be treated as a section in the blog post. 
#     Ensure that every section for the blog has at least 1000 words. 
#     When writing, write like you are explaining to a new coder or someone just starting to learn technology. 
#     Focus on helping folks understand the fundamentals and ensure you do a good job explaining how you move from section of the blog to the next. 
#     Ensure the content flows logically from one section to the next.
#     Use proper Markdown formatting for headers, code blocks, lists, and emphasis where appropriate.
#     """
#     for i, content in enumerate(slides_content, start=1):
#         prompt += f"## Slide {i}\n\n{content}\n\n"
    
#     pieces_client = PiecesClient(config={'baseUrl': 'http://localhost:1000'})
#     response = pieces_client.copilot.question(prompt).answers.iterable[0].text
#     return response

def publish_to_hashnode(title, content, api_key):
    url = "https://api.hashnode.com"
    headers = {
        "Content-Type": "application/json",
        "Authorization": api_key
    }
    query = """
    mutation createPublicationStory($input: CreateStoryInput!) {
        createPublicationStory(input: $input) {
            code
            success
            message
            post {
                _id
                slug
                url
            }
        }
    }
    """
    variables = {
        "input": {
            "title": title,
            "contentMarkdown": content,
            "tags": [],
            "publicationId": "",  # Add your publication ID here
            "isRepublished": {
                "originalArticleURL": ""
            }
        }
    }
    response = requests.post(url, json={"query": query, "variables": variables}, headers=headers)
    return response.json()


# ... (other functions and imports remain the same)


def convert_slides_to_blog_post(slides_content):
    prompt = """You are a helpful assistant that converts slide content into a blog post. 
    Convert the following slide content into a structured blog post using Markdown syntax. 
    Each slide should be treated as a section in the blog post. 
    Ensure that every section for the blog has at least 1000 words. 
    When writing, write like you are explaining to a new coder or someone just starting to learn technology. 
    Focus on helping folks understand the fundamentals and ensure you do a good job explaining how you move from section of the blog to the next. 
    Ensure the content flows logically from one section to the next.
    Use proper Markdown formatting for headers, code blocks, lists, and emphasis where appropriate.
    Also, provide a suitable title for the blog post.
    """
    for i, content in enumerate(slides_content, start=1):
        prompt += f"## Slide {i}\n\n{content}\n\n"
    
    pieces_client = PiecesClient(config={'baseUrl': 'http://localhost:1000'})
    response = pieces_client.copilot.question(prompt).answers.iterable[0].text
    
    # Extract title and content from the response
    lines = response.split('\n')
    title = lines[0].strip('# ')  # Assuming the first line is the title
    content = '\n'.join(lines[1:])
    
    return {"title": title, "content": content}

def main():
    st.title("PowerPoint to Blog Post Converter")

    uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")

    if uploaded_file is not None:
        # Save the uploaded file temporarily
        with open("temp.pptx", "wb") as f:
            f.write(uploaded_file.getvalue())

        # Extract content from the PowerPoint file
        slides_content = extract_text_from_pptx("temp.pptx")

        # Convert slides to blog post
        blog_post = convert_slides_to_blog_post(slides_content)

        # Display the generated blog post
        st.subheader("Generated Blog Post")
        st.markdown(blog_post["content"])

        # Hashnode publishing section
        st.subheader("Publish to Hashnode")
        
        # Use session state to keep track of the publishing process
        if 'publishing_state' not in st.session_state:
            st.session_state.publishing_state = 'initial'

        if st.session_state.publishing_state == 'initial':
            if st.button("Publish to Hashnode"):
                st.session_state.publishing_state = 'api_key_entry'
                st.experimental_rerun()

        if st.session_state.publishing_state == 'api_key_entry':
            api_key = st.text_input("Enter your Hashnode API key", type="password")
            if st.button("Confirm and Publish"):
                if api_key:
                    with st.spinner("Publishing to Hashnode..."):
                        response = publish_to_hashnode(blog_post["title"], blog_post["content"], api_key)
                    if response.get('data', {}).get('createPublicationStory', {}).get('success'):
                        st.success(f"Published successfully! URL: {response['data']['createPublicationStory']['post']['url']}")
                    else:
                        st.error("Failed to publish. Please check your API key and try again.")
                else:
                    st.warning("Please enter your Hashnode API key.")

if __name__ == "__main__":
    main()

# def main():
#     st.title("PowerPoint to Blog Post Converter")

#     # File uploader
#     uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")

#     if uploaded_file is not None:
#         # Save the uploaded file temporarily
#         with open("temp.pptx", "wb") as f:
#             f.write(uploaded_file.getvalue())

#         # Extract content from the PowerPoint file
#         slides_content = extract_text_from_pptx("temp.pptx")

#         # Convert slides to blog post
#         blog_post = convert_slides_to_blog_post(slides_content)

#         # Display the generated blog post
#         st.subheader("Generated Blog Post")
#         st.markdown(blog_post)

#         # Hashnode publishing section
#         st.subheader("Publish to Hashnode")
        
#         # Use session state to keep track of the publishing process
#         if 'publishing_state' not in st.session_state:
#             st.session_state.publishing_state = 'initial'

#         if st.session_state.publishing_state == 'initial':
#             if st.button("Publish to Hashnode"):
#                 st.session_state.publishing_state = 'api_key_entry'
#                 st.experimental_rerun()

#         if st.session_state.publishing_state == 'api_key_entry':
#             api_key = st.text_input("Enter your Hashnode API key", type="password")
#             if st.button("Confirm and Publish"):
#                 if api_key:
#                     with st.spinner("Publishing to Hashnode..."):
#                         response = publish_to_hashnode(blog_post['title'], blog_post['content'], api_key)
#                     if response.get('data', {}).get('createPublicationStory', {}).get('success'):
#                         st.success(f"Published successfully! URL: {response['data']['createPublicationStory']['post']['url']}")
#                     else:
#                         st.error("Failed to publish. Please check your API key and try again.")
#                 else:
#                     st.warning("Please enter your Hashnode API key.")

# if __name__ == "__main__":
#     main()

# def main():
#     st.title("PowerPoint to Blog Post Converter")

#     uploaded_file = st.file_uploader("Choose a PowerPoint file", type="pptx")
    
#     if uploaded_file is not None:
#         slides_content = extract_text_from_pptx(uploaded_file)
#         blog_post = create_blog_post_structure(slides_content)

#         if st.button("Generate Blog Post"):
#             with st.spinner("Generating blog post..."):
#                 ai_generated_post = convert_slides_to_blog_post(slides_content)

#             st.success("Blog post generated successfully!")
#             st.markdown(ai_generated_post)

#             # Option to publish to Hashnode
#             if st.button("Publish to Hashnode"):
#                 api_key = st.text_input("Enter your Hashnode API key", type="password")
#                 if api_key:
#                     with st.spinner("Publishing to Hashnode..."):
#                         response = publish_to_hashnode(blog_post['title'], ai_generated_post, api_key)
#                     if response.get('data', {}).get('createPublicationStory', {}).get('success'):
#                         st.success(f"Published successfully! URL: {response['data']['createPublicationStory']['post']['url']}")
#                     else:
#                         st.error("Failed to publish. Please check your API key and try again.")

# if __name__ == "__main__":
#     main()
